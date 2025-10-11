from fastapi import FastAPI, HTTPException, File, UploadFile, Depends
from fastapi.middleware.cors import CORSMiddleware
from sqlalchemy.orm import Session
from sqlalchemy import text
from app.db import SessionLocal, init_db, StartupDocument, Base, engine
import pdfplumber
from pptx import Presentation
import io
import os
import json
import re
from dotenv import load_dotenv
import google.generativeai as genai
from google.generativeai.types import HarmCategory, HarmBlockThreshold

load_dotenv()

# Verify API key is loaded (but don't print it!)
if not os.getenv("GEMINI_API_KEY"):
    print("⚠️  WARNING: GEMINI_API_KEY not found in environment")
else:
    print("✅ Gemini API key loaded")

app = FastAPI(title="VCMinds Backend API")

app.add_middleware(
    CORSMiddleware,
    allow_origins=[
        "http://localhost:3000", 
        "http://localhost:5173",
        "http://127.0.0.1:3000",
        "http://127.0.0.1:5173"
    ],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Initialize database with error handling
print("=" * 50)
print("DATABASE INITIALIZATION:")
try:
    init_db()
    print("✅ init_db() completed")
except Exception as e:
    print(f"❌ init_db() failed: {e}")

# Debug: Print database connection info
print("DATABASE CONNECTION INFO:")
print(f"Using database: {os.getenv('DATABASE_URL', 'sqlite:///./vcminds.db')[:50]}...")
try:
    with engine.connect() as conn:
        result = conn.execute(text("SELECT 1"))
        print("✅ Database connection successful!")
        
        # Check if table exists
        result = conn.execute(text("""
            SELECT table_name 
            FROM information_schema.tables 
            WHERE table_schema = 'public' AND table_name = 'startup_documents'
        """))
        table_exists = result.fetchone()
        
        if table_exists:
            print("✅ Table 'startup_documents' exists!")
        else:
            print("⚠️  Table 'startup_documents' NOT found - creating manually...")
            # Force create table
            Base.metadata.create_all(bind=engine)
            print("✅ Table created manually!")
            
except Exception as e:
    print(f"❌ Database setup failed: {e}")
print("=" * 50)

def get_db():
    db = SessionLocal()
    try:
        yield db
    finally:
        db.close()

@app.post("/upload_deck/")
async def upload_deck(file: UploadFile = File(...), db: Session = Depends(get_db)):
    filename = file.filename.lower()
    contents = await file.read()
    all_text = ""
    
    if filename.endswith(".pdf"):
        with pdfplumber.open(io.BytesIO(contents)) as pdf:
            all_text = "".join([page.extract_text() or "" for page in pdf.pages])
        doc_type = "pdf"
    elif filename.endswith(".pptx"):
        prs = Presentation(io.BytesIO(contents))
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
        all_text = "\n".join(text_runs)
        doc_type = "pptx"
    else:
        raise HTTPException(status_code=400, detail="Only PDF and PPTX files are supported.")

    doc = StartupDocument(
        file_name=file.filename,
        text=all_text,
        analysis_status="uploaded"
    )
    
    db.add(doc)
    db.commit()
    db.refresh(doc)
    return {"doc_id": doc.id, "file_name": doc.file_name, "doc_type": doc_type}

@app.get("/doc/{doc_id}")
def get_document(doc_id: int, db: Session = Depends(get_db)):
    doc = db.query(StartupDocument).filter(StartupDocument.id == doc_id).first()
    if not doc:
        raise HTTPException(status_code=404, detail="Document not found")
    return {
        "doc_id": doc.id,
        "file_name": doc.file_name,
        "upload_time": doc.upload_time,
        "text": doc.text,
        "analysis_status": doc.analysis_status,
        "analysis_result": doc.analysis_result
    }

@app.post("/analyze/{doc_id}")
def analyze_document(doc_id: int, db: Session = Depends(get_db)):
    doc = db.query(StartupDocument).filter(StartupDocument.id == doc_id).first()
    if not doc:
        raise HTTPException(status_code=404, detail="Document not found")

    pitch_text = doc.text or ""
    
    prompt = f"""
You are a venture capital analyst specializing in startup due diligence.
Analyze the following pitch deck carefully and extract key investment information.
Focus on clarity, accuracy, and investor-relevant insight.

Return ONLY a valid JSON object – no extra text, markdown, or explanations.

**CRITICAL INSTRUCTIONS FOR FOUNDER DETAILS:**
1. Extract ALL mentions of founders, co-founders, team members, and executives
2. For each founder mentioned by name, create a separate entry in "Founder Details"
3. If specific details aren't provided, infer reasonable information from context:
   - Educational background: Look for university names, degrees, certifications
   - Previous experience: Extract company names, roles, industries mentioned
   - Years of experience: Calculate from dates or infer from seniority level
   - Key achievements: Extract ANY accomplishments, awards, or metrics mentioned
4. If a field truly has zero information in the deck, use "Not mentioned in pitch deck" instead of "Not provided"
5. Parse collective team statements (e.g., "+80 years experience") and try to distribute across known founders
6. Look for LinkedIn URLs, email signatures, or social media handles

**EXAMPLE OF GOOD EXTRACTION:**
If deck says "Anthony Lesoisier, CSO. Team has 80+ years investment experience, managed $2B+ assets"
Then extract:
- name: "Anthony Lesoisier"
- title: "Co-founder & CSO"
- previous_experience: "Part of team with 80+ years collective investment experience, $2B+ assets managed"
- years_of_experience: "15-20 years (estimated from team collective experience)"

**IMPORTANT: For 'Overall Score', give an integer from 1 to 10 (NOT 0-100). Score reflects investment quality.**

The response should exactly follow this format:
{{
  "Company Name": "",
  "Founder(s)": "",
  "Problem Statement": "",
  "Solution Overview": "",
  "Market Size": "",
  "Business Model": "",
  "Traction": "",
  "Funding Ask": "",
  "Key Risks/Red Flags": "",
  "1-Sentence Investment Summary": "",
  "Overall Score": 0,
  "Founder Details": [
    {{
      "Founder Name": "",
      "Title": "",
      "Educational Background": "",
      "Previous Experience": "",
      "Previous Startups": "",
      "Key Achievements": "",
      "LinkedIn Mentioned": "",
      "Years of Experience": ""
    }}
  ],
  "Team Composition": {{
    "total_team_size": "",
    "key_roles_covered": "",
    "missing_expertise": "",
    "team_strength_assessment": ""
  }},
  "Competitive Advantages": "",
  "Technology Stack": "",
  "Go-to-Market Strategy": "",
  "Financial Projections": "",
  "Investment Timeline": "",
  "Use of Funds": "",
  "Exit Strategy": ""
}}

Pitch deck text:
{pitch_text[:8000]}
"""

    gemini_api_key = os.getenv("GEMINI_API_KEY")
    if not gemini_api_key:
        raise HTTPException(status_code=500, detail="Gemini API key missing")

    try:
        genai.configure(api_key=gemini_api_key)
        model = genai.GenerativeModel('models/gemini-2.5-flash')
        
        safety_settings = {
            HarmCategory.HARM_CATEGORY_HARASSMENT: HarmBlockThreshold.BLOCK_NONE,
            HarmCategory.HARM_CATEGORY_HATE_SPEECH: HarmBlockThreshold.BLOCK_NONE,
            HarmCategory.HARM_CATEGORY_SEXUALLY_EXPLICIT: HarmBlockThreshold.BLOCK_NONE,
            HarmCategory.HARM_CATEGORY_DANGEROUS_CONTENT: HarmBlockThreshold.BLOCK_NONE,
        }

        response = model.generate_content(prompt, safety_settings=safety_settings)
        raw_content = response.text
        print("RAW CONTENT FROM GEMINI:\n", raw_content[:500])  # Only print first 500 chars

        # Remove code fences
        raw_content = re.sub(r"```json", "", raw_content)
        raw_content = re.sub(r"```", "", raw_content)
        raw_content = raw_content.strip()

        # Extract JSON
        match = re.search(r"\{.*\}", raw_content, re.DOTALL)
        if not match:
            raise HTTPException(status_code=500, detail="No JSON found in model output")
        
        json_str = match.group(0)
        analysis_result = json.loads(json_str)

        # Normalize score to 0-10 scale if it's 0-100
        if "Overall Score" in analysis_result:
            score = analysis_result["Overall Score"]
            if score > 10:
                analysis_result["Overall Score"] = round(score / 10, 1)

        doc.analysis_result = json.dumps(analysis_result)
        doc.analysis_status = "analyzed"
        db.commit()
        db.refresh(doc)
        return {"doc_id": doc.id, "analysis_result": analysis_result}

    except json.JSONDecodeError as e:
        print(f"JSON DECODE ERROR: {e}")
        doc.analysis_status = "error"
        db.commit()
        raise HTTPException(status_code=500, detail=f"JSON parsing failed: {str(e)}")
    except Exception as e:
        print(f"ANALYSIS ERROR: {e}")
        doc.analysis_status = "error"
        db.commit()
        raise HTTPException(status_code=500, detail=f"Analysis failed: {str(e)}")

@app.get("/docs/")
def list_all_docs(db: Session = Depends(get_db)):
    docs = db.query(StartupDocument).filter(StartupDocument.analysis_status == "analyzed").all()
    return {
        "startups": [
            {
                "id": doc.id,
                "file_name": doc.file_name,
                "upload_time": doc.upload_time.isoformat() if doc.upload_time else None,
                "analysis_result": doc.analysis_result
            }
            for doc in docs
        ]
    }

@app.get("/debug/db-status")
def check_database_status():
    """Temporary endpoint to verify database connection"""
    try:
        with engine.connect() as conn:
            # Check if table exists
            result = conn.execute(text("""
                SELECT table_name 
                FROM information_schema.tables 
                WHERE table_schema = 'public'
            """))
            tables = [row[0] for row in result]
            
            return {
                "status": "connected",
                "database_url": os.getenv("DATABASE_URL", "sqlite:///./vcminds.db")[:50] + "...",
                "tables": tables,
                "startup_documents_exists": "startup_documents" in tables
            }
    except Exception as e:
        return {
            "status": "error",
            "error": str(e),
            "database_url": os.getenv("DATABASE_URL", "sqlite:///./vcminds.db")[:50] + "..."
        }

@app.post("/debug/force-create-table")
def force_create_table():
    """Force create the startup_documents table"""
    try:
        Base.metadata.create_all(bind=engine)
        return {"status": "success", "message": "Table creation attempted"}
    except Exception as e:
        return {"status": "error", "error": str(e)}